"""
Gera apresentacao THAUMA para Heads in Health.
python-pptx required.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === CONSTANTES ===
AZUL = RGBColor(0x00, 0x10, 0x70)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CIANO = RGBColor(0x40, 0xD7, 0xFF)
AZUL_CLARO_BG = RGBColor(0xE8, 0xF4, 0xFD)
CINZA_TEXTO = RGBColor(0x33, 0x33, 0x33)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# === HELPERS ===
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_name="Arial",
                font_size=16, bold=False, color=CINZA_TEXTO, alignment=PP_ALIGN.LEFT,
                word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_bullets(slide, left, top, width, height, items, font_size=14,
                color=CINZA_TEXTO, bold_prefix=False, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        if bold_prefix and " — " in item:
            parts = item.split(" — ", 1)
            run1 = p.add_run()
            run1.text = parts[0] + " — "
            run1.font.name = "Arial"
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.name = "Arial"
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
        else:
            p.text = item
            p.font.name = "Arial"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color

        p.level = 0
    return tf

def add_accent_bar(slide, position="top"):
    if position == "top":
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), W, Inches(0.08)
        )
        bar = slide.shapes[-1]
    else:
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), H - Inches(0.08), W, Inches(0.08)
        )
        bar = slide.shapes[-1]
    bar.fill.solid()
    bar.fill.fore_color.rgb = CIANO
    bar.line.fill.background()
    return bar

def add_title_slide(slide, title_text, font_size=32):
    add_accent_bar(slide, "top")
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                title_text, font_size=font_size, bold=True, color=AZUL)

def add_box(slide, left, top, width, height, text, font_size=13,
            fill_color=AZUL_CLARO_BG, border_color=CIANO, text_color=CINZA_TEXTO,
            bold=False, alignment=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    if " — " in text and not bold:
        parts = text.split(" — ", 1)
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = parts[0]
        r1.font.name = "Arial"
        r1.font.size = Pt(font_size)
        r1.font.bold = True
        r1.font.color.rgb = text_color
        r2 = p.add_run()
        r2.text = " — " + parts[1]
        r2.font.name = "Arial"
        r2.font.size = Pt(font_size)
        r2.font.bold = False
        r2.font.color.rgb = text_color
    else:
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color
    p.alignment = alignment
    return shape

def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CIANO
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1 — CAPA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, AZUL)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
            "THAUMA", font_size=72, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
            "Inteligencia & Narrativa em Saude", font_size=28, bold=False,
            color=CIANO, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
            "Apresentacao de Produtos e Capacidades", font_size=16, bold=False,
            color=BRANCO, alignment=PP_ALIGN.CENTER)

# thin ciano line separator
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.3), Inches(5.333), Inches(0.04))
sep.fill.solid()
sep.fill.fore_color.rgb = CIANO
sep.line.fill.background()


# ============================================================
# SLIDE 2 — QUEM SOMOS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "Quem somos")

add_bullets(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5), [
    "Consultoria que transforma dados publicos de saude em inteligencia acionavel",
    "Fundada em dezembro de 2025",
    "Primeiro cliente entregue e validado (Santa Casa de Ouro Preto)",
    "Stack: DATASUS, TSE, BigQuery, IA generativa",
    "Equipe multidisciplinar: dados, marketing, juridico, financas",
], font_size=18, spacing=Pt(12))

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 3 — FILOSOFIA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "De Doxa a Episteme")

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.7),
            "Do amadorismo ao conhecimento verdadeiro",
            font_size=22, bold=True, color=AZUL, alignment=PP_ALIGN.LEFT)

add_bullets(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(3.5), [
    "ALETHEIA  --  Verdade: dados rastreaveis, nunca maquiados",
    "LOGOS  --  Razao: logica estruturada, premissas verificaveis",
    "TECHNE  --  Tecnica: IA e dados como ferramentas de excelencia",
    "PRAXIS  --  Acao: conhecimento so serve se gerar resultado",
], font_size=16, bold_prefix=False, spacing=Pt(10))

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
            "\"O espanto da descoberta. A ciencia do resultado.\"",
            font_size=14, bold=True, color=CIANO, alignment=PP_ALIGN.LEFT)

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 4 — O PROBLEMA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "O Problema")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
            "Hospitais filantropicos dependem de emendas parlamentares para investir e sobreviver.",
            font_size=20, bold=True, color=AZUL, alignment=PP_ALIGN.LEFT)

add_bullets(slide, Inches(0.8), Inches(2.6), Inches(11), Inches(4), [
    "Hoje abordam parlamentares sem dados -- no escuro",
    "Nao sabem quais parlamentares tem eleitores atendidos pelo hospital",
    "Taxa de sucesso baixa, esforco desperdicado",
    "Resultado: milhoes em emendas que nunca chegam",
], font_size=16, spacing=Pt(10))

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 5 — A SOLUCAO: PRISMA DE CAPTACAO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "Prisma de Captacao")

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
            "O Kit de Emendas Parlamentares", font_size=20, bold=False, color=CIANO)

add_textbox(slide, Inches(0.8), Inches(1.9), Inches(11), Inches(0.5),
            "Cruzamento inedito entre DATASUS e TSE", font_size=16, color=CINZA_TEXTO)

# 3 boxes
box_w = Inches(3.5)
box_h = Inches(2.5)
y = Inches(3.0)

add_box(slide, Inches(0.8), y, box_w, box_h,
        "DATASUS\n\nQuem o hospital atende:\nmunicipios, procedimentos, perfil de pacientes",
        font_size=14, bold=True)

add_box(slide, Inches(4.9), y, box_w, box_h,
        "TSE\n\nOnde o parlamentar tem votos:\nbase eleitoral por municipio",
        font_size=14, bold=True)

add_box(slide, Inches(9.0), y, box_w, box_h,
        "RANKING\n\nQuais parlamentares abordar com maior chance de sucesso",
        font_size=14, bold=True, fill_color=CIANO, text_color=BRANCO, border_color=AZUL)

# arrows between boxes
add_arrow(slide, Inches(4.35), Inches(3.9), Inches(0.5), Inches(0.35))
add_arrow(slide, Inches(8.45), Inches(3.9), Inches(0.5), Inches(0.35))

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 6 — SCORE SAT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "Score SAT -- Alinhamento Territorial")

add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
            "SAT = (Votos / 1.000)  x  (Pacientes / 100)",
            font_size=32, bold=True, color=AZUL, alignment=PP_ALIGN.CENTER)

# separator
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2.9), Inches(7), Inches(0.03))
sep.fill.solid()
sep.fill.fore_color.rgb = CIANO
sep.line.fill.background()

add_bullets(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(3), [
    "Votos do parlamentar no municipio",
    "Pacientes do municipio atendidos pelo hospital",
    "Quanto maior o SAT, maior o alinhamento de interesses",
], font_size=18, spacing=Pt(12))

add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
            "Metodologia proprietaria THAUMA", font_size=14, bold=True,
            color=CIANO, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 7 — OS 4 COMPONENTES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "4 Componentes do Prisma")

components = [
    "Dossie de Evidencias — Perfil completo do hospital: municipios, procedimentos, CIDs, demografia",
    "Radar Politico — Ranking de parlamentares por Score SAT",
    "Dialetica de Convencimento — Argumentos personalizados para cada parlamentar-alvo",
    "Retorica da Influencia — Textos prontos: pitch, justificativa tecnica, mensagens para redes sociais",
]

bw = Inches(5.5)
bh = Inches(1.2)
x_left = Inches(0.8)
x_right = Inches(6.8)
y1 = Inches(1.8)
y2 = Inches(3.5)

positions = [(x_left, y1), (x_right, y1), (x_left, y2), (x_right, y2)]

for idx, (comp, (x, y)) in enumerate(zip(components, positions)):
    num = str(idx + 1)
    # number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = AZUL
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.name = "Arial"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BRANCO
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_box(slide, x + Inches(0.55), y, bw - Inches(0.55), bh, comp, font_size=13)

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 8 — FORMATOS DE ENTREGA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "O que o hospital recebe")

items = [
    "Dashboard HTML interativo (auto-contido, abre em qualquer navegador)",
    "Planilha Excel formatada (9 abas com dados e analises)",
    "Roteiros de abordagem parlamentar",
    "Textos prontos para comunicacao",
    "Apresentacoes de apoio",
]

add_bullets(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5), items,
            font_size=18, spacing=Pt(14))

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 9 — CASE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "Case Validado")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
            "Santa Casa de Misericordia de Ouro Preto",
            font_size=24, bold=True, color=AZUL)

add_bullets(slide, Inches(0.8), Inches(2.4), Inches(11), Inches(3.5), [
    "Primeiro cliente -- contrato assinado e executado",
    "Entrega completa em abril de 2026",
    "Materiais validados pela diretoria do hospital",
    "Dashboard interativo + Excel + roteiros + textos",
], font_size=16, spacing=Pt(12))

add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
            "Referencia disponivel mediante solicitacao",
            font_size=14, bold=True, color=CIANO)

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 10 — TRANSICAO ASSISTENCIAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "Alem da Captacao Politica")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1.0),
            "A mesma base de dados que identifica oportunidades politicas tambem revela oportunidades assistenciais.",
            font_size=18, color=CINZA_TEXTO)

# visual: two boxes with arrow
box_left_w = Inches(4.5)
box_h2 = Inches(2.0)
y_boxes = Inches(3.5)

add_box(slide, Inches(1.2), y_boxes, box_left_w, box_h2,
        "Inteligencia Politica\n\nCaptacao de emendas,\nalinhamento parlamentar,\nScore SAT",
        font_size=15, bold=True, fill_color=AZUL_CLARO_BG, border_color=CIANO)

add_arrow(slide, Inches(5.9), Inches(4.2), Inches(1.2), Inches(0.5))

add_box(slide, Inches(7.4), y_boxes, box_left_w, box_h2,
        "Inteligencia Assistencial\n\nVazios assistenciais,\nfaturamento SUS,\ncapacidade diagnostica",
        font_size=15, bold=True, fill_color=CIANO, text_color=BRANCO, border_color=AZUL)

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 11 — FONTES DE DADOS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "A Base: DATASUS")

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
            "A maior base de dados de saude publica do mundo",
            font_size=18, bold=False, color=CIANO)

add_bullets(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(4), [
    "SIH — Sistema de Informacoes Hospitalares (internacoes, AIH)",
    "SIA — Sistema de Informacoes Ambulatoriais (BPA, APAC)",
    "CNES — Cadastro Nacional de Estabelecimentos de Saude",
    "SIGTAP — Tabela de Procedimentos SUS",
    "TSE — Tribunal Superior Eleitoral (dados eleitorais)",
    "IBGE — Dados socioeconomicos e geograficos",
], font_size=16, bold_prefix=True, spacing=Pt(10))

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
            "Data Lake proprio no BigQuery -- +1.6 milhao de registros (MG 2025)",
            font_size=14, bold=True, color=CIANO)

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 12 — CAPACIDADES ASSISTENCIAIS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "Inteligencia Assistencial")

cap_items = [
    "Vazio Assistencial — Municipios desatendidos por especialidade ou procedimento. Onde estao as oportunidades de expansao?",
    "Radar SUS / Faturamento — Monitoramento continuo da producao hospitalar. Rejeicoes, glosas, mix de procedimentos.",
    "Capacidade Diagnostica — Gestao de equipamentos e imagens. Onde ha ociosidade ou sobrecarga?",
]

bw3 = Inches(11.5)
bh3 = Inches(1.3)
y_start = Inches(1.8)

for i, item in enumerate(cap_items):
    add_box(slide, Inches(0.8), y_start + Inches(i * 1.6), bw3, bh3, item, font_size=15)

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 13 — PIPELINE DE DADOS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BRANCO)
add_title_slide(slide, "Como Funciona")

pipeline_steps = [
    "Dados Publicos\n(DATASUS/FTP)",
    "ETL Automatizado\n(R + microdatasus)",
    "Data Lake\n(BigQuery)",
    "Analises Customizadas\n(Python + SQL)",
    "Entregaveis Visuais\n(Dashboards, Relatorios)",
]

box_w_pipe = Inches(2.1)
box_h_pipe = Inches(1.6)
y_pipe = Inches(3.0)
x_start = Inches(0.5)
gap = Inches(0.45)

for i, step in enumerate(pipeline_steps):
    x = x_start + i * (box_w_pipe + gap)
    if i == len(pipeline_steps) - 1:
        add_box(slide, x, y_pipe, box_w_pipe, box_h_pipe, step, font_size=13,
                bold=True, fill_color=CIANO, text_color=BRANCO, border_color=AZUL,
                alignment=PP_ALIGN.CENTER)
    else:
        add_box(slide, x, y_pipe, box_w_pipe, box_h_pipe, step, font_size=13,
                bold=True, alignment=PP_ALIGN.CENTER)
    # arrow after box (except last)
    if i < len(pipeline_steps) - 1:
        arrow_x = x + box_w_pipe + Inches(0.05)
        add_arrow(slide, arrow_x, y_pipe + Inches(0.6), Inches(0.35), Inches(0.3))

add_accent_bar(slide, "bottom")


# ============================================================
# SLIDE 14 — FECHAMENTO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, AZUL)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
            "Dados que provocam espanto.",
            font_size=40, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(1.0),
            "Estrategias que geram resultado.",
            font_size=40, bold=True, color=CIANO, alignment=PP_ALIGN.CENTER)

# separator
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5), Inches(5.333), Inches(0.03))
sep.fill.solid()
sep.fill.fore_color.rgb = CIANO
sep.line.fill.background()

add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "Pedro William Ribeiro Diniz",
            font_size=18, bold=True, color=BRANCO, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
            "pedro@thaumaconsultoria.com.br",
            font_size=16, color=CIANO, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.4),
            "THAUMA -- Inteligencia & Narrativa em Saude",
            font_size=14, color=BRANCO, alignment=PP_ALIGN.CENTER)


# === SALVAR ===
output_path = r"C:\Users\User\thauma-staff\Projetos\Heads_in_Health\Apresentacao_THAUMA_Produtos.pptx"
prs.save(output_path)
print(f"Apresentacao salva em: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
